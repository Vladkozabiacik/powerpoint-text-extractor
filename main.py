import os
from pptx import Presentation
from docx import Document

def extract_text_from_shape(shape):
    if hasattr(shape, "text_frame"):
        return shape.text_frame.text
    else:
        return ""

def process_pptx_files(directory):
    for root, dirs, files in os.walk(directory):
        doc = Document()
        for file in files:
            if file.endswith(".pptx"):
                pptx_file = os.path.join(root, file)
                print(f"Extracting text from {pptx_file}...")
                doc.add_heading(f"From presentation -> {os.path.basename(pptx_file)}:", level=1)
                prs = Presentation(pptx_file)
                empty_line = False
                for slide in prs.slides:
                    for shape in slide.shapes:
                        text = extract_text_from_shape(shape)
                        filtered_text = "".join(c for c in text if c != '\x0b')
                        if any(c.strip() for c in filtered_text):
                            if empty_line:
                                doc.add_paragraph()
                                empty_line = False
                            doc.add_paragraph(filtered_text)
                        else:
                            empty_line = True
                output_file = os.path.join(root, "output.docx")
                doc.save(output_file)
                print(f"Text extracted and saved to {output_file}")

root_directory = os.getcwd()
process_pptx_files(root_directory)
